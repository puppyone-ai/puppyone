"""
本地文件解析工具
使用 Python 库解析 PDF、DOCX、PPTX 和图片文件解析为 Markdown
"""
from pathlib import Path
from typing import Optional
from app.utils.logger import log_info, log_error


class LocalFileParser:
    """本地文件解析器"""
    
    def __init__(self):
        pass
    
    def parse_pdf(self, file_path: Path) -> str:
        """
        解析 PDF 文件为 Markdown
        
        Args:
            file_path: PDF 文件路径
        
        Returns:
            Markdown 格式的文本内容
        """
        try:
            # 尝试使用 pdfplumber（更好的表格和布局支持）
            try:
                import pdfplumber
                return self._parse_pdf_with_pdfplumber(file_path, pdfplumber)
            except ImportError:
                # 回退到 PyPDF2
                import PyPDF2
                return self._parse_pdf_with_pypdf2(file_path, PyPDF2)
        except ImportError:
            raise ImportError("需要安装 pdfplumber 或 PyPDF2: pip install pdfplumber 或 pip install PyPDF2")
        except Exception as e:
            log_error(f"Failed to parse PDF: {e}")
            raise Exception(f"PDF 解析失败: {str(e)}")
    
    def _parse_pdf_with_pdfplumber(self, file_path: Path, pdfplumber) -> str:
        """使用 pdfplumber 解析 PDF"""
        markdown_parts = []
        
        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            log_info(f"Parsing PDF with pdfplumber: {total_pages} pages")
            
            for page_num, page in enumerate(pdf.pages, 1):
                # 提取文本
                text = page.extract_text()
                if text:
                    markdown_parts.append(f"## 第 {page_num} 页\n\n{text}\n")
                
                # 提取表格
                tables = page.extract_tables()
                if tables:
                    for table_num, table in enumerate(tables, 1):
                        markdown_parts.append(f"\n### 表格 {table_num}\n\n")
                        markdown_parts.append(self._table_to_markdown(table))
                        markdown_parts.append("\n")
        
        return "\n".join(markdown_parts)
    
    def _parse_pdf_with_pypdf2(self, file_path: Path, PyPDF2) -> str:
        """使用 PyPDF2 解析 PDF（回退方案）"""
        markdown_parts = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            total_pages = len(pdf_reader.pages)
            log_info(f"Parsing PDF with PyPDF2: {total_pages} pages")
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text:
                    markdown_parts.append(f"## 第 {page_num} 页\n\n{text}\n")
        
        return "\n".join(markdown_parts)
    
    def parse_docx(self, file_path: Path) -> str:
        """
        解析 DOCX 文件为 Markdown
        
        Args:
            file_path: DOCX 文件路径
        
        Returns:
            Markdown 格式的文本内容
        """
        try:
            from docx import Document
        except ImportError:
            raise ImportError("需要安装 python-docx: pip install python-docx")
        
        try:
            doc = Document(file_path)
            markdown_parts = []
            
            log_info(f"Parsing DOCX: {file_path.name}")
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                # 根据样式判断标题级别
                style_name = para.style.name.lower()
                if 'heading 1' in style_name or 'title' in style_name:
                    markdown_parts.append(f"# {text}\n")
                elif 'heading 2' in style_name:
                    markdown_parts.append(f"## {text}\n")
                elif 'heading 3' in style_name:
                    markdown_parts.append(f"### {text}\n")
                elif 'heading 4' in style_name:
                    markdown_parts.append(f"#### {text}\n")
                else:
                    markdown_parts.append(f"{text}\n")
            
            # 提取表格
            for table in doc.tables:
                markdown_parts.append("\n")
                markdown_parts.append(self._docx_table_to_markdown(table))
                markdown_parts.append("\n")
            
            return "\n".join(markdown_parts)
        except Exception as e:
            log_error(f"Failed to parse DOCX: {e}")
            raise Exception(f"DOCX 解析失败: {str(e)}")
    
    def parse_pptx(self, file_path: Path) -> str:
        """
        解析 PPTX 文件为 Markdown
        
        Args:
            file_path: PPTX 文件路径
        
        Returns:
            Markdown 格式的文本内容
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("需要安装 python-pptx: pip install python-pptx")
        
        try:
            prs = Presentation(file_path)
            markdown_parts = []
            
            log_info(f"Parsing PPTX: {file_path.name}, {len(prs.slides)} slides")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                markdown_parts.append(f"## 幻灯片 {slide_num}\n\n")
                
                # 提取幻灯片中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # 检查是否是标题
                        if shape == slide.shapes.title:
                            markdown_parts.append(f"### {text}\n\n")
                        else:
                            markdown_parts.append(f"{text}\n\n")
                
                markdown_parts.append("\n")
            
            return "\n".join(markdown_parts)
        except Exception as e:
            log_error(f"Failed to parse PPTX: {e}")
            raise Exception(f"PPTX 解析失败: {str(e)}")
    
    def parse_image(self, file_path: Path) -> str:
        """
        解析图片文件为 Markdown（提取元数据）
        
        Args:
            file_path: 图片文件路径
        
        Returns:
            Markdown 格式的元数据信息
        """
        try:
            from PIL import Image
            from PIL.ExifTags import TAGS
        except ImportError:
            raise ImportError("需要安装 Pillow: pip install Pillow")
        
        try:
            with Image.open(file_path) as img:
                width, height = img.size
                format_name = img.format
                mode = img.mode
                file_size = file_path.stat().st_size
                
                markdown_parts = [
                    f"# 图片信息\n\n",
                    f"**文件名**: {file_path.name}\n",
                    f"**尺寸**: {width} × {height} 像素\n",
                    f"**格式**: {format_name}\n",
                    f"**颜色模式**: {mode}\n",
                    f"**文件大小**: {file_size / 1024:.2f} KB\n",
                ]
                
                # 提取 EXIF 信息
                try:
                    exifdata = img.getexif()
                    if exifdata:
                        markdown_parts.append("\n## EXIF 信息\n\n")
                        for tag_id in exifdata:
                            tag = TAGS.get(tag_id, tag_id)
                            data = exifdata.get(tag_id)
                            if data:
                                markdown_parts.append(f"- **{tag}**: {data}\n")
                except Exception as e:
                    log_info(f"Could not extract EXIF data: {e}")
                
                markdown_parts.append("\n\n> 注意：图片内容需要使用 OCR 工具提取文本")
                
                return "".join(markdown_parts)
        except Exception as e:
            log_error(f"Failed to parse image: {e}")
            raise Exception(f"图片解析失败: {str(e)}")
    
    def parse_file(self, file_path: Path, file_extension: str) -> str:
        """
        根据文件扩展名解析文件
        
        Args:
            file_path: 文件路径
            file_extension: 文件扩展名（小写，不含点）
        
        Returns:
            Markdown 格式的文本内容
        """
        extension = file_extension.lower()
        
        if extension == 'pdf':
            return self.parse_pdf(file_path)
        elif extension in ['docx']:
            return self.parse_docx(file_path)
        elif extension in ['pptx']:
            return self.parse_pptx(file_path)
        elif extension in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif']:
            return self.parse_image(file_path)
        else:
            raise ValueError(f"不支持的文件格式: {extension}")
    
    def _table_to_markdown(self, table: list) -> str:
        """将表格转换为 Markdown 格式"""
        if not table or not table[0]:
            return ""
        
        markdown_lines = []
        
        # 表头
        header = table[0]
        markdown_lines.append("| " + " | ".join(str(cell) if cell else "" for cell in header) + " |")
        markdown_lines.append("| " + " | ".join("---" for _ in header) + " |")
        
        # 数据行
        for row in table[1:]:
            markdown_lines.append("| " + " | ".join(str(cell) if cell else "" for cell in row) + " |")
        
        return "\n".join(markdown_lines)
    
    def _docx_table_to_markdown(self, table) -> str:
        """将 DOCX 表格转换为 Markdown 格式"""
        if not table.rows:
            return ""
        
        markdown_lines = []
        
        # 表头
        header_row = table.rows[0]
        header_cells = [cell.text.strip() for cell in header_row.cells]
        markdown_lines.append("| " + " | ".join(header_cells) + " |")
        markdown_lines.append("| " + " | ".join("---" for _ in header_cells) + " |")
        
        # 数据行
        for row in table.rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            markdown_lines.append("| " + " | ".join(cells) + " |")
        
        return "\n".join(markdown_lines)


def get_local_file_parser() -> LocalFileParser:
    """获取本地文件解析器实例"""
    return LocalFileParser()

